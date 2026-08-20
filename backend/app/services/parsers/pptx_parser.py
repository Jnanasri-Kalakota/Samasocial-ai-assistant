import io
from typing import List
from pptx import Presentation
from app.services.parsers.base import BaseParser, ExtractedContent
from app.core.exceptions import DocumentParsingError
import logging

logger = logging.getLogger(__name__)

class PPTXParser(BaseParser):
    async def parse(self, file_bytes: bytes, filename: str = "presentation.pptx", **kwargs) -> List[ExtractedContent]:
        results: List[ExtractedContent] = []
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            total_slides = len(prs.slides)
            
            for idx, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                slide_title = None
                
                # Check for slide title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = slide.shapes.title.text.strip()
                    slide_texts.append(f"Title: {slide_title}")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        continue
                    if hasattr(shape, "text") and shape.text:
                        clean_shape_text = shape.text.strip()
                        if clean_shape_text:
                            slide_texts.append(clean_shape_text)
                            
                combined_text = "\n".join(slide_texts).strip()
                if not combined_text:
                    continue
                    
                label = f"Slide {idx}"
                if slide_title:
                    label = f"Slide {idx}: {slide_title[:30]}"
                    
                results.append(
                    ExtractedContent(
                        text=combined_text,
                        source_type="pptx",
                        source_name=filename,
                        location_label=label,
                        location_meta={
                            "slide_number": idx,
                            "slide_title": slide_title,
                            "total_slides": total_slides,
                            "filename": filename
                        }
                    )
                )
        except Exception as e:
            logger.error(f"Failed to parse PPTX {filename}: {e}")
            raise DocumentParsingError(f"Error parsing PowerPoint '{filename}': {str(e)}")
            
        if not results:
            raise DocumentParsingError(f"No text content found in PowerPoint '{filename}'")
            
        return results
